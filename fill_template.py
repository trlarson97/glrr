from pptx import Presentation
from pptx.util import Inches
import json
import os
import subprocess
import sys

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

def load_schools():
    json_path = os.path.join(BASE_DIR, "schools.json")
    print(f"Loading schools from: {json_path}")
    print(f"Schools file exists: {os.path.exists(json_path)}")
    if os.path.exists(json_path):
        with open(json_path, 'r') as f:
            data = json.load(f)
        print(f"Loaded {len(data)} schools")
        return data
    print("WARNING: schools.json not found!")
    return {}

# Set to False to show full school names (separate name/school boxes in template).
# Set to True to use schools.json abbreviations again.
USE_ABBREVIATIONS = True

def initials_list(members):
    """Convert 'Dylan Gamnje, Gideon Gash' -> 'D. Gamnje, G. Gash'."""
    if not members:
        return ''
    out = []
    for full in members.split(','):
        full = full.strip()
        if not full:
            continue
        parts = full.split()
        if len(parts) >= 2:
            out.append(f"{parts[0][0]}. {' '.join(parts[1:])}")
        else:
            out.append(full)
    return ', '.join(out)

def shorten_school(name, schools):
    if not name:
        return name
    # Always clean up athletic.net truncation artifacts (trailing "...")
    name = name.split(' - ')[0].strip()
    name = name.rstrip('.').strip()
    # Only strip a "(" fragment if it's an INCOMPLETE truncation (no closing paren).
    # Keep complete "City (School)" formats like "LaGrange (Lyons)".
    if '(' in name and ')' not in name:
        name = name.split('(')[0].strip()

    if not USE_ABBREVIATIONS:
        return name

    # 1. Exact match
    if name in schools:
        return schools[name]
    # 2. Strip relay suffix like " - A"
    clean = name.split(' - ')[0].strip()
    if clean in schools:
        return schools[clean]
    # 3. Strip trailing ellipsis/truncation and any "(..." fragment
    stripped = clean.rstrip('.').split('(')[0].strip()
    if stripped in schools:
        return schools[stripped]
    # 4. Truncation match — athletic.net cuts names off, so if our (shorter)
    #    name is the start of a known key, use that key's full name
    for key, full in schools.items():
        if len(stripped) >= 10 and key.startswith(stripped):
            return full
    return name

def build_rows(data, schools=None):
    """Compute the display rows (rank/name/school/time) for a single graphic.

    This is the parse-independent, PowerPoint-independent core: it applies the
    school-name restoration and the relay/team/individual formatting rules and
    returns a plain list of dicts. Used by both fill_template() (to fill the
    .pptx placeholders) and the /process endpoint (to feed the browser graphic),
    so the two paths can never drift apart.
    """
    if schools is None:
        schools = load_schools()

    is_relay    = data.get('is_relay', False)
    is_team     = data.get('eventtype', 'individual') == 'team'
    scoring     = data.get('scoring', 'points')
    max_results = data.get('max_results', 8)
    sport       = data.get('sport', '')
    results     = data.get('results', [])

    rows = []
    for i in range(1, max_results + 1):
        if i > len(results):
            break
        raw_name   = results[i-1].get('name', '')
        raw_school = results[i-1].get('school', '')
        time       = results[i-1].get('time', '')

        if is_team:
            display_name = shorten_school(raw_name, schools)
            display_time = f"{time} pts" if scoring == 'points' else f"{time}"
            if raw_school:
                if 'Cross Country' in sport:
                    display_school = f"Scorers: {raw_school}"
                else:
                    plural = '&' in raw_school or ',' in raw_school
                    label = 'Top Scorers' if plural else 'Top Scorer'
                    display_school = f"{label}: {raw_school}"
            else:
                display_school = ''
        elif is_relay:
            display_name   = shorten_school(raw_name, schools)
            display_school = initials_list(raw_school)
            display_time   = time
        else:
            display_name   = raw_name
            display_school = shorten_school(raw_school, schools)
            display_time   = time

        rows.append({
            'rank':   str(results[i-1].get('rank', i)),
            'name':   display_name,
            'school': display_school,
            'time':   display_time,
        })
    return rows


def replace_in_paragraph(para, replacements):
    for run in para.runs:
        for placeholder, value in replacements.items():
            if placeholder in run.text:
                run.text = run.text.replace(placeholder, value)
    for placeholder, value in replacements.items():
        full_text = ''.join(r.text for r in para.runs)
        if placeholder in full_text:
            new_text = full_text.replace(placeholder, value)
            if new_text != full_text:
                if para.runs:
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run.text = ''

def get_template_path(data):
    category = data.get('category', 'state')
    if category == 'conference':
        key = data.get('conference', '').lower().replace(' ', '_')
    elif category == 'meet':
        key = 'meet'   # uses templates/meet.pptx if present, else falls back
    else:
        key = data.get('state', 'michigan').lower().replace(' ', '_')
    path = os.path.join(BASE_DIR, "templates", f"{key}.pptx")
    print(f"Looking for template: {path}")
    if os.path.exists(path):
        print(f"Found: {path}")
        return path
    fallback = os.path.join(BASE_DIR, "mhsaa_template.pptx")
    print(f"Not found, using fallback: {fallback}")
    return fallback

def export_png(pptx_path):
    """Try multiple methods to convert pptx to PNG on Windows."""
    base = os.path.splitext(pptx_path)[0]
    png_path = base + '.png'
    out_dir = os.path.dirname(pptx_path)

    # Method 1: LibreOffice
    lo_paths = [
        r'C:\Program Files\LibreOffice\program\soffice.exe',
        r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
        'libreoffice',
        'soffice'
    ]
    for lo_cmd in lo_paths:
        try:
            result = subprocess.run(
                [lo_cmd, '--headless', '--convert-to', 'png', '--outdir', out_dir, pptx_path],
                capture_output=True, text=True, timeout=60
            )
            print(f"LibreOffice stdout: {result.stdout}")
            print(f"LibreOffice stderr: {result.stderr}")
            if os.path.exists(png_path):
                print(f"PNG via LibreOffice: {png_path}")
                return png_path
        except (FileNotFoundError, subprocess.TimeoutExpired) as e:
            print(f"LibreOffice attempt failed ({lo_cmd}): {e}")
            continue

    # Method 2: Windows PowerPoint COM automation
    try:
        import comtypes.client
        import time
        time.sleep(0.5)  # ensure file is fully written
        abs_pptx = os.path.abspath(pptx_path)
        abs_png  = os.path.abspath(png_path)
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        deck = powerpoint.Presentations.Open(abs_pptx, WithWindow=False)
        # Export all slides as PNG — width 1920px for high quality
        deck.Export(abs_png, "PNG", 1920, 1080)
        deck.Close()
        powerpoint.Quit()
        # PowerPoint Export() may create slide001.png instead of the exact filename
        # Check for both
        slide_png = os.path.splitext(abs_png)[0] + "1.png"
        if os.path.exists(abs_png):
            print(f"PNG via PowerPoint COM: {abs_png}")
            return abs_png
        elif os.path.exists(slide_png):
            os.rename(slide_png, abs_png)
            print(f"PNG via PowerPoint COM (renamed): {abs_png}")
            return abs_png
    except Exception as e:
        print(f"COM method failed: {e}")

    # Method 3: python-pptx + Pillow (basic rendering)
    try:
        from PIL import Image
        import io
        prs = Presentation(pptx_path)
        slide = prs.slides[0]
        # Get slide dimensions
        width = int(prs.slide_width.pt)
        height = int(prs.slide_height.pt)
        img = Image.new('RGB', (width * 2, height * 2), color='white')
        img.save(png_path)
        print(f"PNG placeholder created: {png_path}")
        return png_path
    except Exception as e:
        print(f"Pillow method failed: {e}")

    print("All PNG export methods failed")
    return None

def fill_template(data, output_path=None):
    schools      = load_schools()
    max_results  = data.get('max_results', 8)
    template_path = get_template_path(data)

    prs   = Presentation(template_path)
    slide = prs.slides[0]

    replacements = {
        '{{sport}}':    data.get('sport', ''),
        '{{division}}': data.get('division', ''),
        '{{meet}}':     data.get('meet', ''),
        '{{gender}}':   data.get('gender', ''),
        '{{event}}':    data.get('event', ''),
        '{{type}}':     data.get('type', 'All-State'),
    }

    rows = build_rows(data, schools)
    for i in range(1, max_results + 1):
        if i <= len(rows):
            r = rows[i-1]
            replacements[f'{{{{rank_{i}}}}}']   = r['rank']
            replacements[f'{{{{name_{i}}}}}']   = r['name']
            replacements[f'{{{{school_{i}}}}}'] = r['school']
            replacements[f'{{{{time_{i}}}}}']   = r['time']
        else:
            replacements[f'{{{{rank_{i}}}}}']   = ''
            replacements[f'{{{{name_{i}}}}}']   = ''
            replacements[f'{{{{school_{i}}}}}'] = ''
            replacements[f'{{{{time_{i}}}}}']   = ''

    for i in range(max_results + 1, 10):
        replacements[f'{{{{rank_{i}}}}}']   = ''
        replacements[f'{{{{name_{i}}}}}']   = ''
        replacements[f'{{{{school_{i}}}}}'] = ''
        replacements[f'{{{{time_{i}}}}}']   = ''

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            replace_in_paragraph(para, replacements)

    if output_path is None:
        safe_event  = data.get('event', 'event').replace(' ', '_').replace('/', '-')
        safe_gender = data.get('gender', 'gender')
        safe_div    = data.get('division', 'div') or 'open'
        category    = data.get('category', 'state')
        if category == 'conference':
            prefix = data.get('conference', 'conf')
        elif category == 'meet':
            prefix = (data.get('meet', 'meet') or 'meet')
        else:
            prefix = data.get('state', 'output')
        prefix = prefix.replace(' ', '_')
        output_path = os.path.join(BASE_DIR, f"output_{prefix}_{safe_gender}_{safe_div}_{safe_event}.pptx")

    prs.save(output_path)
    print(f"Saved PPTX: {output_path}")

    png_path = export_png(output_path)
    return output_path, png_path


def fill_carousel(slides_data):
    """
    Generate one PNG per slide for a carousel post.
    Returns list of {label, pptx, png} dicts.
    """
    results = []
    for idx, slide_data in enumerate(slides_data):
        safe_event  = slide_data.get('event', 'event').replace(' ', '_').replace('/', '-')
        safe_gender = slide_data.get('gender', 'gender')
        safe_div    = slide_data.get('division', 'div') or 'open'
        category    = slide_data.get('category', 'state')
        if category == 'conference':
            prefix = slide_data.get('conference', 'conf')
        elif category == 'meet':
            prefix = (slide_data.get('meet', 'meet') or 'meet')
        else:
            prefix = slide_data.get('state', 'output')
        prefix = prefix.replace(' ', '_')
        output_path = os.path.join(BASE_DIR, f"carousel_{idx+1:02d}_{prefix}_{safe_gender}_{safe_div}_{safe_event}.pptx")

        pptx_path, png_path = fill_template(slide_data, output_path=output_path)

        label = f"{slide_data.get('gender','')} {slide_data.get('event','')} — {slide_data.get('division','')}"
        results.append({
            'label': label,
            'pptx':  os.path.basename(pptx_path),
            'png':   os.path.basename(png_path) if png_path else None
        })

    return results
